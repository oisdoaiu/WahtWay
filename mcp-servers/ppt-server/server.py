"""
PPT MCP Server — 通过 python-pptx 操作 PowerPoint 文件
暴露 3 个 Tool：read-ppt-structure / clone-and-fill / create-ppt
"""
import json
import os
import copy
import tempfile
from pathlib import Path

from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import Tool, TextContent

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════════
# LAYOUT DETECTION
# ═══════════════════════════════════════════════════════════════

def detect_layout_type(slide, idx, total):
    """检测单张幻灯片的布局类型"""
    shapes = slide.shapes
    texts = []
    for shape in shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)

    has_image = any(hasattr(s, 'image') for s in shapes)
    has_table = bool(slide.shapes._spTree.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tbl'))

    # 封面特征：文本少、标题短、第一页
    if idx == 0 and len(texts) <= 4 and texts[0] if texts else False:
        if len(texts[0]) < 60:
            return "cover"

    # 图表
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/chart}'
    has_chart = bool(slide.shapes._spTree.findall(f'.//{ns}chart'))
    if has_chart:
        return "chart"

    if has_table:
        return "table"

    # 多卡片（3-4 个大小相似的形状 + 短语）
    if len(texts) >= 3:
        short_count = sum(1 for t in texts if len(t) < 30)
        if short_count >= 3:
            return "three-cards"

    # 大数字
    if texts and texts[0].replace('%', '').replace(',', '').replace('.', '').isdigit():
        return "big-number"

    # 流程/时间轴（多个短文本）
    if 3 <= len(texts) <= 8 and all(len(t) < 40 for t in texts):
        return "process"

    if idx == total - 1 and len(texts) <= 3:
        return "closing"

    return "content"


def read_ppt_structure(filepath):
    """读取 PPTX 结构，返回 JSON"""
    if not os.path.exists(filepath):
        return f"文件不存在: {filepath}"

    try:
        prs = Presentation(filepath)
        slides = []
        total = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)

            layout_type = detect_layout_type(slide, i, total)
            features = {}
            if layout_type == "chart":
                features["hasChart"] = True
            if layout_type == "table":
                features["hasTable"] = True

            slides.append({
                "index": i + 1,
                "title": texts[0] if texts else "",
                "content": texts[1:],
                "textCount": len(texts),
                "layoutType": layout_type,
                "features": features,
            })

        return json.dumps({
            "file": os.path.basename(filepath),
            "totalSlides": total,
            "slides": slides,
        }, ensure_ascii=False, indent=2)

    except Exception as e:
        return f"读取 PPT 失败: {e}"


# ═══════════════════════════════════════════════════════════════
# SLIDE CLONING
# ═══════════════════════════════════════════════════════════════

def _shape_yx(shape):
    """返回形状的 (top, left) 坐标用于排序"""
    try:
        return (shape.top, shape.left)
    except:
        return (0, 0)


def _clone_slide(prs, source_slide):
    """深拷贝一张幻灯片，复用空白版式避免重复 layout"""
    # 固定使用空白版式（index 6），不引入新的 slideLayout
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    # 复制源幻灯片的所有形状
    for shape in source_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # 复制关系（图片等外部资源）
    for rel in source_slide.part.rels.values():
        if rel.is_external:
            new_slide.part.rels.get_or_add(rel.reltype, rel.target)

    # 删除由版式带来的占位符
    for ph in list(new_slide.placeholders):
        sp = ph._element
        if sp.getparent() is not None:
            sp.getparent().remove(sp)

    return new_slide


def _fill_slide_text(slide, title=None, bullets=None):
    """向幻灯片填充文字：按视觉位置从上到下、从左到右替换可编辑文本"""
    texts_to_fill = []
    if title:
        texts_to_fill.append(title)
    if bullets:
        texts_to_fill.extend(bullets)

    if not texts_to_fill:
        return

    # 收集所有 run，按位置排序
    runs_with_pos = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        pos = _shape_yx(shape)
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    runs_with_pos.append((pos, run))

    runs_with_pos.sort(key=lambda x: (x[0][0], x[0][1]))  # top first, then left

    for ti in range(min(len(texts_to_fill), len(runs_with_pos))):
        _, run = runs_with_pos[ti]
        run.text = texts_to_fill[ti]


def clone_and_fill(template_path, slide_map, output_path):
    """从模板克隆指定页、填充内容、输出新 PPTX"""
    if not os.path.exists(template_path):
        return f"模板文件不存在: {template_path}"

    try:
        prs = Presentation(template_path)
        total = len(prs.slides)

        # 创建空白输出
        new_prs = Presentation()
        # 设置 16:9
        new_prs.slide_width = prs.slide_width
        new_prs.slide_height = prs.slide_height

        for item in slide_map:
            tpl_idx = int(item.get("templateSlide", 1)) - 1
            if tpl_idx < 0 or tpl_idx >= total:
                return f"模板幻灯片序号无效: {tpl_idx + 1}（模板共 {total} 页）"

            source_slide = prs.slides[tpl_idx]
            new_slide = _clone_slide(new_prs, source_slide)
            _fill_slide_text(new_slide, item.get("title"), item.get("bullets"))

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        new_prs.save(output_path)

        return json.dumps({
            "status": "success",
            "output": output_path,
            "slidesGenerated": len(slide_map),
            "templateUsed": os.path.basename(template_path),
            "message": f"从模板 {os.path.basename(template_path)} 克隆了 {len(slide_map)} 页 → {output_path}",
        }, ensure_ascii=False)

    except Exception as e:
        return f"模板克隆失败: {e}"


# ═══════════════════════════════════════════════════════════════
# CREATE PPT (from scratch)
# ═══════════════════════════════════════════════════════════════

THEMES = {
    "business":  {"bg": "FFFFFF", "title": "1A3C6E", "text": "333333", "accent": "2E75B6"},
    "modern":    {"bg": "FFFFFF", "title": "2D3436", "text": "636E72", "accent": "0984E3"},
    "warm":      {"bg": "FFF8F0", "title": "C0392B", "text": "555555", "accent": "E67E22"},
    "dark":      {"bg": "2C3E50", "title": "ECF0F1", "text": "BDC3C7", "accent": "3498DB"},
    "minimal":   {"bg": "FFFFFF", "title": "111111", "text": "666666", "accent": "999999"},
    "creative":  {"bg": "FFFFFF", "title": "6C5CE7", "text": "444444", "accent": "FD79A8"},
}

COLORS = ["2E75B6", "E17055", "00B894", "6C5CE7", "0984E3", "FD79A8"]


def _add_top_line(slide, accent, w=10):
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(w), Inches(0.05)).line.fill.background()


def _hex(hex_str):
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def create_ppt(slides_data, output_path, theme_name="business"):
    """从结构化数据生成 PPTX"""
    theme = THEMES.get(theme_name, THEMES["business"])

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]  # blank layout

        for i, slide_data in enumerate(slides_data):
            slide = prs.slides.add_slide(blank)
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = _hex(theme["bg"])

            layout = slide_data.get("layout", "content")
            title = slide_data.get("title", "")
            bullets = slide_data.get("bullets", [])

            # 顶部装饰线
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = _hex(theme["accent"])
            line.line.fill.background()

            if layout == "cover":
                _render_cover(slide, slide_data, theme)
            elif layout == "closing":
                _render_closing(slide, slide_data, theme)
            elif layout == "three-cards":
                _render_cards(slide, slide_data, theme)
            elif layout == "big-number":
                _render_big_number(slide, slide_data, theme)
            elif layout == "process":
                _render_process(slide, slide_data, theme)
            else:
                _render_content(slide, slide_data, theme, i)

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)

        return json.dumps({
            "status": "success",
            "output": output_path,
            "slidesGenerated": len(slides_data),
            "theme": theme_name,
            "message": f"创建了 {len(slides_data)} 页 PPT → {output_path}（主题: {theme_name}）",
        }, ensure_ascii=False)

    except Exception as e:
        return f"创建 PPT 失败: {e}"


def _add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color="333333", align="left", font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _hex(color)
    p.font.name = font_name
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    return txBox


def _render_cover(slide, data, theme):
    # 左侧大色块
    left_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5), Inches(7.5))
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = _hex(theme["accent"])
    left_block.line.fill.background()

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    author = data.get("author", "")

    if title:
        _add_text_box(slide, 5.5, 2.0, 7.0, 1.5, title, 42, True, theme["title"])
    if subtitle:
        _add_text_box(slide, 5.5, 3.6, 7.0, 0.7, subtitle, 18, False, theme["text"])
    if author:
        _add_text_box(slide, 5.5, 4.4, 7.0, 0.5, author, 13, False, "999999")


def _render_closing(slide, data, theme):
    thanks = data.get("title", "谢谢")
    _add_text_box(slide, 0, 2.5, 13.3, 1.5, thanks, 40, True, theme["accent"], "center")

    bullets = data.get("bullets", [])
    if bullets:
        text = " · ".join(bullets)
        _add_text_box(slide, 2, 4.2, 9.3, 1.0, text, 16, False, theme["text"], "center")


def _render_cards(slide, data, theme):
    title = data.get("title", "")
    if title:
        _add_text_box(slide, 1.0, 0.4, 11.3, 0.7, title, 26, True, theme["title"], "center")

    cards = data.get("cards", [])
    from pptx.enum.shapes import MSO_SHAPE
    card_w, gap = 3.8, 0.4

    for ci in range(min(len(cards), 3)):
        card = cards[ci]
        cx = 0.7 + ci * (card_w + gap)

        # 卡片背景
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(1.5), Inches(card_w), Inches(4.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex("FFFFFF")
        rect.line.color.rgb = _hex("E0E0E0")
        rect.line.width = Pt(0.5)

        # 数字
        num_color = COLORS[ci % len(COLORS)]
        _add_text_box(slide, cx + 0.5, 1.8, card_w - 1.0, 1.0, card.get("num", ""), 36, True, num_color, "center")

        # 标题
        card_title = card.get("title", "")
        _add_text_box(slide, cx + 0.3, 3.0, card_w - 0.6, 0.5, card_title, 16, True, theme["title"], "center")

        # 描述
        card_desc = card.get("desc", "")
        _add_text_box(slide, cx + 0.3, 3.6, card_w - 0.6, 2.0, card_desc, 11, False, theme["text"], "center")


def _render_big_number(slide, data, theme):
    num = data.get("num", "")
    title = data.get("title", "")
    bullets = data.get("bullets", [])

    # 背景圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1.2), Inches(5.3), Inches(4.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = _hex(theme["accent"])
    circle.line.fill.background()
    # 设置透明度（需要 lxml）
    try:
        from lxml import etree
        solidFill = circle.fill._fill
        srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', '15000')  # 15% opacity = 85% transparency
    except:
        pass

    _add_text_box(slide, 4.5, 2.0, 4.3, 1.8, num, 60, True, theme["title"], "center")
    _add_text_box(slide, 2, 5.0, 9.3, 0.6, title, 20, True, theme["title"], "center")
    if bullets:
        _add_text_box(slide, 3, 5.6, 7.3, 0.5, bullets[0], 13, False, theme["text"], "center")


def _render_process(slide, data, theme):
    title = data.get("title", "")
    if title:
        _add_text_box(slide, 1.0, 0.4, 11.3, 0.7, title, 26, True, theme["title"], "center")

    steps = data.get("steps", data.get("bullets", []))
    step_count = min(len(steps), 5)
    step_w, gap = 2.0, 0.6
    total_w = step_count * step_w + (step_count - 1) * gap
    start_x = (13.3 - total_w) / 2

    for pi in range(step_count):
        px = start_x + pi * (step_w + gap)
        color = COLORS[pi % len(COLORS)]

        # 圆角矩形
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(2.0), Inches(step_w), Inches(1.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex(color)
        rect.line.fill.background()

        _add_text_box(slide, px + 0.1, 2.1, step_w - 0.2, 1.3, steps[pi], 13, True, "FFFFFF", "center")

        # 编号圆
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px + step_w/2 - 0.25), Inches(3.8), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex(color)
        circle.line.fill.background()
        _add_text_box(slide, px + step_w/2 - 0.25, 3.85, 0.5, 0.4, str(pi + 1), 12, True, "FFFFFF", "center")

        # 箭头连接
        if pi < step_count - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(px + step_w), Inches(2.5), Inches(gap), Inches(0.5))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = _hex(theme["accent"])
            arrow.line.fill.background()


def _render_content(slide, data, theme, page_idx):
    title = data.get("title", "")
    bullets = data.get("bullets", [])

    if title:
        _add_text_box(slide, 0.8, 0.5, 11.7, 0.8, title, 28, True, theme["title"])
        # 标题下划线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = _hex(theme["accent"])
        line.line.fill.background()

    if bullets:
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.3), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for bi, bullet in enumerate(bullets):
            if bi == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.font.color.rgb = _hex(theme["text"])
            p.font.name = "Microsoft YaHei"
            p.space_after = Pt(8)
            p.level = 0

    # 页码
    _add_text_box(slide, 12.2, 6.8, 0.8, 0.4, str(page_idx + 1), 10, False, theme["accent"], "right")


# ═══════════════════════════════════════════════════════════════
# MCP SERVER
# ═══════════════════════════════════════════════════════════════

server = Server("ppt-mcp-server")


@server.list_tools()
async def list_tools():
    return [
        Tool(
            name="read-ppt-structure",
            description="读取 PPTX 文件的结构和内容，返回每页的文本、布局类型，自动识别封面/内容/图表/卡片/流程/大数字等页面类型。用于分析模板、提取内容。",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {"type": "string", "description": "PPTX 文件的完整路径"},
                },
                "required": ["filepath"],
            },
        ),
        Tool(
            name="clone-and-fill",
            description="从模板 PPTX 中克隆指定页面，填充新内容，生成新 PPT。核心能力：可以选择模板的任意页面（可重复选择同一页），克隆其完整视觉设计（颜色/字体/背景/形状/Logo），只替换文本。这是套用模板的正确方式——选页、克隆、填充。",
            inputSchema={
                "type": "object",
                "properties": {
                    "templatePath": {"type": "string", "description": "模板 PPTX 文件路径"},
                    "slideMap": {
                        "type": "array",
                        "description": "幻灯片映射。每项 { templateSlide: 数字(从1开始，选择模板中的第几页), title: 标题文本, bullets: [要点数组] }。相同 templateSlide 可重复出现来克隆同一页多次。",
                        "items": {
                            "type": "object",
                            "properties": {
                                "templateSlide": {"type": "number", "description": "模板中的页码，从 1 开始"},
                                "title": {"type": "string", "description": "替换后的标题"},
                                "bullets": {"type": "array", "items": {"type": "string"}, "description": "替换后的要点列表"},
                            },
                        },
                    },
                    "outputPath": {"type": "string", "description": "输出 .pptx 文件路径"},
                },
                "required": ["templatePath", "slideMap", "outputPath"],
            },
        ),
        Tool(
            name="create-ppt",
            description="从零创建 PowerPoint 演示文稿，提供结构化幻灯片数据。支持 6 种内置主题（business/modern/warm/dark/minimal/creative），支持多种布局（cover/content/closing/three-cards/big-number/process）。用于没有模板时从文案生成视觉化 PPT。",
            inputSchema={
                "type": "object",
                "properties": {
                    "slides": {
                        "type": "array",
                        "description": "幻灯片数组。每项 { title, bullets, layout?: cover|content|closing|three-cards|big-number|process, subtitle?, author?, cards?: [{num,title,desc}], num?, steps? }",
                    },
                    "outputPath": {"type": "string", "description": "输出 .pptx 文件路径"},
                    "theme": {"type": "string", "description": "主题: business/modern/warm/dark/minimal/creative（默认 business）"},
                },
                "required": ["slides", "outputPath"],
            },
        ),
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict):
    if name == "read-ppt-structure":
        result = read_ppt_structure(str(arguments["filepath"]))
    elif name == "clone-and-fill":
        result = clone_and_fill(
            str(arguments["templatePath"]),
            list(arguments["slideMap"]),
            str(arguments["outputPath"]),
        )
    elif name == "create-ppt":
        result = create_ppt(
            list(arguments["slides"]),
            str(arguments["outputPath"]),
            str(arguments.get("theme", "business")),
        )
    else:
        result = f"未知工具: {name}"

    return [TextContent(type="text", text=result)]


async def main():
    async with stdio_server() as (read_stream, write_stream):
        await server.run(read_stream, write_stream, server.create_initialization_options())


if __name__ == "__main__":
    import asyncio
    asyncio.run(main())
